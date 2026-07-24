"""PPTX renderer — lossless write-back of translations into the original file.

The renderer reopens the *original* pptx (guaranteeing nothing else is lost),
locates each shape/paragraph/cell via the block's SourceRef, and replaces the
text while preserving the first run's formatting. Localized images are swapped
in by replacing the underlying image part's blob.
"""

from __future__ import annotations

import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from translayer.fonts.layout import fit_text
from translayer.fonts.registry import FontRegistry
from translayer.ir.models import Block, DocumentIR
from translayer.languages import pptx_language_tag
from translayer.plugins import registry

# --------------------------------------------------------------------------- #
# XML namespaces
# --------------------------------------------------------------------------- #
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qtag(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


class PptxRenderer:
    name = "pptx"

    def supported_formats(self) -> list[str]:
        return ["pptx"]

    def render(self, ir: DocumentIR, input_path: str, output_path: str) -> None:
        prs = Presentation(input_path)

        # Index shapes by (slide_index, shape_id) including nested groups.
        shape_index: dict[tuple[int, int], object] = {}
        visible_text_widths: dict[tuple[int, int], int] = {}
        for s_idx, slide in enumerate(prs.slides):
            shapes = list(slide.shapes)
            for shape_index_on_slide, shape in enumerate(shapes):
                self._index_shape(shape, s_idx, shape_index)
                if getattr(shape, "has_text_frame", False):
                    visible_text_widths[(s_idx, int(shape.shape_id))] = min(
                        self._visible_text_width(
                            shape, shapes[shape_index_on_slide + 1 :]
                        ),
                        self._container_text_width(
                            shape, shapes[:shape_index_on_slide]
                        ),
                    )

        images = {im.id: im for im in ir.resources.images}
        # Cache loaded SmartArt data parts and their XML roots so multiple
        # points on the same diagram share one in-memory tree and one write-back.
        smartart_cache: dict[object, ET.Element] = {}

        for block in ir.blocks:
            if block.type == "image":
                self._render_image(block, shape_index, images)
                continue
            if block.type == "smartart":
                self._render_smartart(
                    block,
                    shape_index,
                    smartart_cache,
                    target_lang=ir.meta.target_lang,
                )
                continue
            if block.target_text is None:
                continue
            if block.target_text == block.source_text:
                continue
            self._render_text_block(
                block,
                shape_index,
                target_lang=ir.meta.target_lang,
                fit_width_emu=visible_text_widths.get(
                    (block.source_ref.slide_index, block.source_ref.shape_id)
                ),
            )

        # Persist any modified SmartArt data parts.
        self._persist_smartart(smartart_cache)

        prs.save(output_path)

    # ------------------------------------------------------------------ #
    def _index_shape(self, shape, slide_index, index) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                self._index_shape(sub, slide_index, index)
            return
        index[(slide_index, int(shape.shape_id))] = shape

    def _render_text_block(
        self,
        block: Block,
        index,
        target_lang: str,
        fit_width_emu: int | None = None,
    ) -> None:
        ref = block.source_ref
        shape = index.get((ref.slide_index, ref.shape_id))
        if shape is None:
            return
        if fit_width_emu and fit_width_emu < int(shape.width):
            shape.width = fit_width_emu

        if ref.kind == "table_cell":
            if not shape.has_table:
                return
            try:
                cell = shape.table.rows[ref.row].cells[ref.col]
            except (IndexError, AttributeError):
                return
            tf = cell.text_frame
        else:
            if not shape.has_text_frame:
                return
            tf = shape.text_frame

        try:
            para = tf.paragraphs[ref.paragraph_index]
        except (IndexError, TypeError):
            return

        self._set_paragraph_text(
            para,
            block.target_text or "",
            language_tag=pptx_language_tag(target_lang),
            target_lang=target_lang,
        )
        self._apply_overflow(
            tf,
            block,
            target_lang,
            fit_width_emu=fit_width_emu,
        )

    @staticmethod
    def _set_paragraph_text(
        para,
        text: str,
        language_tag: str | None = None,
        target_lang: str | None = None,
    ) -> None:
        """Replace paragraph text, preserving the first run's formatting."""
        runs = para.runs
        if not runs:
            run = para.add_run()
            run.text = text
            PptxRenderer._apply_run_language(run, language_tag, target_lang)
            return
        runs[0].text = text
        PptxRenderer._apply_run_language(runs[0], language_tag, target_lang)
        for extra in runs[1:]:
            extra.text = ""

    @staticmethod
    def _apply_run_language(run, language_tag: str | None, target_lang: str | None) -> None:
        if not language_tag:
            return
        properties = run._r.get_or_add_rPr()
        properties.set("lang", language_tag)
        if target_lang != "zh" or properties.find(_qtag(_NS_A, "ea")) is not None:
            return
        east_asian = properties.makeelement(
            _qtag(_NS_A, "ea"), {"typeface": "Microsoft YaHei"}
        )
        properties.append(east_asian)

    @staticmethod
    def _apply_overflow(
        text_frame,
        block: Block,
        target_lang: str,
        *,
        fit_width_emu: int | None = None,
    ) -> None:
        if not block.constraints.can_shrink_font:
            return
        try:
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if not text_frame.text:
                return

            source_font = block.runs[0].font if block.runs else None
            max_size = (
                source_font.size
                if source_font and source_font.size
                else block.layout.base_font.size
                if block.layout and block.layout.base_font.size
                else 32.0
            )
            font_file = FontRegistry().font_for_lang(target_lang)
            position = block.layout.position if block.layout else None
            if position is None or block.source_ref.kind == "table_cell":
                return
            horizontal_margins = int(text_frame.margin_left or 0) + int(
                text_frame.margin_right or 0
            )
            vertical_margins = int(text_frame.margin_top or 0) + int(
                text_frame.margin_bottom or 0
            )
            box_w_px = max(
                1,
                int(
                    round(
                        ((fit_width_emu or position.w) - horizontal_margins)
                        / 12700
                        * 96
                        / 72
                    )
                ),
            )
            box_h_px = max(
                1,
                int(round((position.h - vertical_margins) / 12700 * 96 / 72)),
            )
            max_size_px = max(1, int(round(max_size * 96 / 72)))
            min_size_px = min(
                max_size_px,
                max(4, int(round(block.constraints.min_font_size * 96 / 72))),
            )
            best_size_px, _ = fit_text(
                text_frame.text,
                box_w_px,
                box_h_px,
                font_file,
                max_size=max_size_px,
                min_size=min_size_px,
            )
            best_size = best_size_px * 72 / 96
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        run.font.size = Pt(best_size)
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except (AttributeError, OSError, TypeError, ValueError):
            pass

    @staticmethod
    def _visible_text_width(shape, foreground_shapes) -> int:
        """Return text width before a foreground picture obscures its right side."""
        left = int(shape.left)
        top = int(shape.top)
        right = left + int(shape.width)
        bottom = top + int(shape.height)
        visible_right = right
        for foreground in foreground_shapes:
            if foreground.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            foreground_left = int(foreground.left)
            foreground_top = int(foreground.top)
            foreground_bottom = foreground_top + int(foreground.height)
            vertical_overlap = max(
                0,
                min(bottom, foreground_bottom) - max(top, foreground_top),
            )
            if vertical_overlap < min(int(shape.height), int(foreground.height)) * 0.25:
                continue
            if left < foreground_left < visible_right:
                visible_right = foreground_left
        visible_width = visible_right - left
        return visible_width if visible_width >= int(shape.width) * 0.25 else int(shape.width)

    @staticmethod
    def _container_text_width(shape, background_shapes) -> int:
        """Fit text inside the nearest filled shape that visually contains it."""
        left = int(shape.left)
        top = int(shape.top)
        right = left + int(shape.width)
        bottom = top + int(shape.height)
        for background in reversed(background_shapes):
            if background.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if getattr(background, "has_text_frame", False) and background.text.strip():
                continue
            try:
                if background.fill.type is None:
                    continue
            except (AttributeError, ValueError):
                continue
            background_left = int(background.left)
            background_top = int(background.top)
            background_right = background_left + int(background.width)
            background_bottom = background_top + int(background.height)
            if not (
                background_left <= left
                and background_top <= top
                and background_bottom >= bottom
                and left < background_right < right
            ):
                continue
            left_inset = left - background_left
            inner_right = background_right - left_inset
            container_width = inner_right - left
            if container_width >= int(shape.width) * 0.25:
                return container_width
        return int(shape.width)

    def _render_image(self, block: Block, index, images) -> None:
        ref = block.source_ref
        if ref.image_id is None:
            return
        res = images.get(ref.image_id)
        if res is None or not res.localized_data_ref:
            return
        shape = index.get((ref.slide_index, ref.shape_id))
        if shape is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return
        try:
            with open(res.localized_data_ref, "rb") as fh:
                new_blob = fh.read()
            self._replace_image_blob(shape, new_blob)
            if res.reset_crop_on_render:
                shape.crop_left = 0
                shape.crop_top = 0
                shape.crop_right = 0
                shape.crop_bottom = 0
        except (AttributeError, OSError, KeyError, ValueError):
            pass

    @staticmethod
    def _replace_image_blob(shape, new_blob: bytes) -> None:
        """Replace the underlying image part blob for a picture shape.

        python-pptx 1.0+ exposes |Image| as an immutable value object, so we
        locate the actual |ImagePart| via the shape's blip relationship and
        overwrite its ``_blob``.
        """
        blip = next(shape._element.iter(_qtag(_NS_A, "blip")))
        rId = blip.attrib[_qtag(_NS_R, "embed")]
        image_part = shape.part.related_part(rId)
        image_part._blob = new_blob

    # ------------------------------------------------------------------ #
    # SmartArt (diagram) rendering
    # ------------------------------------------------------------------ #
    def _render_smartart(
        self,
        block: Block,
        index,
        smartart_cache: dict[object, ET.Element],
        target_lang: str,
    ) -> None:
        if block.target_text is None:
            return
        ref = block.source_ref
        if ref.kind != "smartart_point" or not ref.region_id:
            return

        shape = index.get((ref.slide_index, ref.shape_id))
        if shape is None:
            return

        data_part = self._smartart_data_part(shape)
        if data_part is None:
            return

        root = smartart_cache.get(data_part)
        if root is None:
            try:
                root = ET.fromstring(data_part.blob)
                smartart_cache[data_part] = root
            except ET.ParseError:
                return

        model_id = ref.region_id
        for pt in root.findall(f".//{_qtag(_NS_DGM, 'pt')}"):
            if pt.get("modelId") == model_id:
                self._set_point_text(
                    pt,
                    block.target_text or "",
                    language_tag=pptx_language_tag(target_lang),
                    target_lang=target_lang,
                )
                break

    @staticmethod
    def _smartart_data_part(shape):
        relids = shape._element.find(f".//{_qtag(_NS_DGM, 'relIds')}")
        if relids is None:
            return None
        dm_rId = relids.attrib.get(_qtag(_NS_R, "dm"))
        if not dm_rId:
            return None
        return shape.part.related_part(dm_rId)

    @staticmethod
    def _set_point_text(
        pt,
        text: str,
        language_tag: str = "zh-CN",
        target_lang: str = "zh",
    ) -> None:
        """Replace all text inside a diagram point, preserving line breaks as paragraphs."""
        t_elem = pt.find(_qtag(_NS_DGM, "t"))
        if t_elem is None:
            return

        # Keep body properties / list style if present.
        body_pr = t_elem.find(_qtag(_NS_A, "bodyPr"))
        lst_style = t_elem.find(_qtag(_NS_A, "lstStyle"))
        t_elem.clear()
        if body_pr is not None:
            t_elem.append(body_pr)
        if lst_style is not None:
            t_elem.append(lst_style)

        for line in (text or "").split("\n"):
            para = ET.SubElement(t_elem, _qtag(_NS_A, "p"))
            run = ET.SubElement(para, _qtag(_NS_A, "r"))
            rPr = ET.SubElement(run, _qtag(_NS_A, "rPr"))
            rPr.set("lang", language_tag)
            if target_lang == "zh":
                east_asian = ET.SubElement(rPr, _qtag(_NS_A, "ea"))
                east_asian.set("typeface", "Microsoft YaHei")
            text_elem = ET.SubElement(run, _qtag(_NS_A, "t"))
            text_elem.text = line

    @staticmethod
    def _persist_smartart(smartart_cache: dict[object, ET.Element]) -> None:
        """Write modified SmartArt XML back to their package parts.

        Register stable prefixes so PowerPoint recognizes the namespaces and
        the diagram renders correctly.
        """
        ET.register_namespace("dgm", _NS_DGM)
        ET.register_namespace("a", _NS_A)
        ET.register_namespace("r", _NS_R)
        ET.register_namespace("dsp", "http://schemas.microsoft.com/office/drawing/2008/diagram")

        for data_part, root in smartart_cache.items():
            xml_bytes = ET.tostring(root, encoding="UTF-8", xml_declaration=True)
            # ET.tostring does not emit standalone="yes"; PowerPoint expects it
            # on diagram data parts, so patch it in.
            xml_bytes = xml_bytes.replace(
                b'<?xml version=\'1.0\' encoding=\'UTF-8\'?>',
                b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
            )
            data_part._blob = xml_bytes


registry.register("renderer", "pptx")(PptxRenderer)
