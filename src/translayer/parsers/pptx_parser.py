"""PPTX parser — extracts text, layout and images into DocumentIR.

Blocks are paragraph-level (best granularity for translation context and
quality). Each block carries a precise SourceRef so the renderer can write the
translation back losslessly into the original file.
"""

from __future__ import annotations

import os
import tempfile
import xml.etree.ElementTree as ET
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from translayer.ir.models import (
    Block,
    DocMeta,
    DocumentIR,
    Font,
    ImageResource,
    Layout,
    Position,
    Run,
    SourceRef,
)
from translayer.plugins import registry

# --------------------------------------------------------------------------- #
# XML namespaces
# --------------------------------------------------------------------------- #
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qtag(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def _color_hex(font) -> str | None:
    try:
        if font.color and font.color.type is not None and font.color.rgb is not None:
            return f"#{str(font.color.rgb)}"
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _run_font(run) -> Font:
    f = run.font
    return Font(
        name=f.name,
        size=float(f.size.pt) if f.size is not None else None,
        color=_color_hex(f),
        bold=bool(f.bold) if f.bold is not None else False,
        italic=bool(f.italic) if f.italic is not None else False,
        underline=bool(f.underline) if f.underline is not None else False,
    )


def _position(shape) -> Position | None:
    try:
        if None in (shape.left, shape.top, shape.width, shape.height):
            return None
        return Position(
            x=int(shape.left), y=int(shape.top),
            w=int(shape.width), h=int(shape.height),
        )
    except (AttributeError, TypeError):
        return None


def _base_font(shape) -> Font:
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                return _run_font(run)
    except (AttributeError, ValueError):
        pass
    return Font()


def _semantic_role(shape) -> str | None:
    try:
        if shape.is_placeholder:
            idx_type = shape.placeholder_format.type
            name = str(idx_type).lower() if idx_type is not None else ""
            if "title" in name and "subtitle" not in name:
                return "title"
            if "subtitle" in name or "ctr_title" in name:
                return "subtitle"
            if "body" in name:
                return "body"
    except (AttributeError, ValueError):
        pass
    return None


class PptxParser:
    name = "pptx"

    def supported_formats(self) -> list[str]:
        return ["pptx"]

    def parse(self, input_path: str, meta: dict) -> DocumentIR:
        asset_dir = meta.get("asset_dir") or tempfile.mkdtemp(prefix="translayer_assets_")
        os.makedirs(asset_dir, exist_ok=True)

        prs = Presentation(input_path)
        doc = DocumentIR(
            meta=DocMeta(
                source_lang=meta.get("source_lang", "en"),
                target_lang=meta.get("target_lang", "zh"),
                doc_type="pptx",
                title=meta.get("title"),
                glossary_ref=meta.get("glossary_ref"),
                engine_hints={
                    "slide_width": int(prs.slide_width),
                    "slide_height": int(prs.slide_height),
                },
            )
        )

        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                self._walk_shape(shape, slide_index, doc, asset_dir)

        return doc

    # ------------------------------------------------------------------ #
    def _walk_shape(self, shape, slide_index: int, doc: DocumentIR, asset_dir: str) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                self._walk_shape(sub, slide_index, doc, asset_dir)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._parse_picture(shape, slide_index, doc, asset_dir)
            return

        if shape.has_table:
            self._parse_table(shape, slide_index, doc)
            return

        if self._is_smartart(shape):
            self._parse_smartart(shape, slide_index, doc)
            return

        if shape.has_text_frame:
            self._parse_text_frame(shape, slide_index, doc)

    def _parse_text_frame(self, shape, slide_index: int, doc: DocumentIR) -> None:
        role = _semantic_role(shape)
        block_type = role if role in ("title", "subtitle", "body") else "shape_text"
        layout = Layout(
            slide_index=slide_index,
            shape_id=int(shape.shape_id),
            position=_position(shape),
            base_font=_base_font(shape),
        )
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            runs = [Run(text=r.text, font=_run_font(r)) for r in para.runs]
            text = "".join(r.text for r in para.runs)
            if not text.strip():
                continue
            doc.blocks.append(
                Block(
                    id=f"s{slide_index}-sh{shape.shape_id}-p{p_idx}",
                    type=block_type,  # type: ignore[arg-type]
                    semantic_role=role,
                    runs=runs,
                    source_text=text,
                    layout=layout,
                    source_ref=SourceRef(
                        kind="shape_text",
                        slide_index=slide_index,
                        shape_id=int(shape.shape_id),
                        paragraph_index=p_idx,
                    ),
                )
            )

    def _parse_table(self, shape, slide_index: int, doc: DocumentIR) -> None:
        table = shape.table
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                for p_idx, para in enumerate(cell.text_frame.paragraphs):
                    text = "".join(run.text for run in para.runs)
                    if not text.strip():
                        continue
                    runs = [Run(text=run.text, font=_run_font(run)) for run in para.runs]
                    doc.blocks.append(
                        Block(
                            id=f"s{slide_index}-sh{shape.shape_id}-r{r_idx}-c{c_idx}-p{p_idx}",
                            type="table_cell",
                            runs=runs,
                            source_text=text,
                            layout=Layout(
                                slide_index=slide_index,
                                shape_id=int(shape.shape_id),
                                position=_position(shape),
                            ),
                            source_ref=SourceRef(
                                kind="table_cell",
                                slide_index=slide_index,
                                shape_id=int(shape.shape_id),
                                row=r_idx,
                                col=c_idx,
                                paragraph_index=p_idx,
                            ),
                        )
                    )

    def _parse_picture(self, shape, slide_index: int, doc: DocumentIR, asset_dir: str) -> None:
        try:
            image = shape.image
        except (AttributeError, ValueError):
            return
        img_id = f"s{slide_index}-sh{shape.shape_id}-img"
        ext = image.ext or "png"
        out_path = os.path.join(asset_dir, f"{img_id}.{ext}")
        with open(out_path, "wb") as fh:
            fh.write(image.blob)
        px_w, px_h = image.size if image.size else (0, 0)
        doc.resources.images.append(
            ImageResource(
                id=img_id,
                media_type=image.content_type or "image/png",
                data_ref=out_path,
                width=int(px_w),
                height=int(px_h),
            )
        )
        # An image is also a block so the pipeline can address it for rendering.
        doc.blocks.append(
            Block(
                id=img_id,
                type="image",
                translatable=False,
                source_text="",
                layout=Layout(
                    slide_index=slide_index,
                    shape_id=int(shape.shape_id),
                    position=_position(shape),
                ),
                source_ref=SourceRef(
                    kind="image_region",
                    slide_index=slide_index,
                    shape_id=int(shape.shape_id),
                    image_id=img_id,
                ),
            )
        )

    # ------------------------------------------------------------------ #
    # SmartArt (diagram) support
    # ------------------------------------------------------------------ #
    @staticmethod
    def _is_smartart(shape) -> bool:
        """True if shape is a graphicFrame referencing diagram data."""
        elem = getattr(shape, "_element", None)
        if elem is None:
            return False
        return elem.tag == _qtag(_NS_P, "graphicFrame")

    def _parse_smartart(self, shape, slide_index: int, doc: DocumentIR) -> None:
        """Extract translatable text points from a SmartArt diagram."""
        data_part = self._smartart_data_part(shape)
        if data_part is None:
            return

        try:
            root = ET.fromstring(data_part.blob)
        except ET.ParseError:
            return

        layout = Layout(
            slide_index=slide_index,
            shape_id=int(shape.shape_id),
            position=_position(shape),
        )

        for p_idx, pt in enumerate(root.findall(f".//{_qtag(_NS_DGM, 'pt')}")):
            text = self._point_text(pt)
            if not text.strip():
                continue
            model_id = pt.get("modelId", f"pt{p_idx}")
            doc.blocks.append(
                Block(
                    id=f"s{slide_index}-sh{shape.shape_id}-sa{model_id}",
                    type="smartart",
                    source_text=text,
                    layout=layout,
                    source_ref=SourceRef(
                        kind="smartart_point",
                        slide_index=slide_index,
                        shape_id=int(shape.shape_id),
                        paragraph_index=p_idx,
                        region_id=model_id,
                    ),
                )
            )

    @staticmethod
    def _smartart_data_part(shape) -> Any | None:
        """Return the diagram data Part for a SmartArt graphicFrame shape."""
        relids = shape._element.find(f".//{_qtag(_NS_DGM, 'relIds')}")
        if relids is None:
            return None
        dm_rId = relids.attrib.get(_qtag(_NS_R, "dm"))
        if not dm_rId:
            return None
        try:
            return shape.part.related_part(dm_rId)
        except (KeyError, ValueError):
            return None

    @staticmethod
    def _point_text(pt) -> str:
        """Concatenate all <a:t> text inside a diagram point."""
        return "".join(
            t.text or "" for t in pt.findall(f".//{_qtag(_NS_A, 't')}")
        )


registry.register("parser", "pptx")(PptxParser)


# Keep linter happy about the imported alias used for EMU typing clarity.
_ = Emu
