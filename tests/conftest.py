"""Shared test fixtures."""

from __future__ import annotations

import os

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def sample_pptx(tmp_path) -> str:
    """A small deck: title+body slide, a table, and a picture with text."""
    prs = Presentation()

    # Slide 1: title + body
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Results"
    body = slide.placeholders[1]
    body.text_frame.text = "Revenue grew strongly."
    p = body.text_frame.add_paragraph()
    p.text = "Costs were controlled."

    # Slide 2: table + picture
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    rows, cols = 2, 2
    table = slide2.shapes.add_table(
        rows, cols, Inches(0.5), Inches(0.5), Inches(4), Inches(1)
    ).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Users"
    table.cell(1, 1).text = "1000"

    # Build a picture containing text
    img_path = os.path.join(tmp_path, "chart.png")
    img = Image.new("RGB", (400, 200), (235, 240, 250))
    ImageDraw.Draw(img).text((20, 80), "Revenue", fill=(0, 0, 0))
    img.save(img_path)
    slide2.shapes.add_picture(img_path, Inches(1), Inches(2), Inches(4), Inches(2))

    out = os.path.join(tmp_path, "sample.pptx")
    prs.save(out)
    return out
