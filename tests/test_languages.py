from __future__ import annotations

import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from translayer.languages import (
    language_name,
    normalize_language,
    paddle_language,
    pptx_language_tag,
    tesseract_language,
)
from translayer.localize.text_pipeline import localize_text
from translayer.pipeline import parse_document, render_document

PHRASES = {
    "en": ("Quarterly Results", "Revenue increased significantly."),
    "de": ("Quartalsergebnisse", "Der Umsatz ist deutlich gestiegen."),
    "zh": ("季度业绩", "收入显著增长。"),
}

TRANSLATIONS = {
    (source, target, source_text): target_text
    for source, source_phrases in PHRASES.items()
    for target, target_phrases in PHRASES.items()
    if source != target
    for source_text, target_text in zip(source_phrases, target_phrases, strict=True)
}


class _MatrixTranslator:
    def translate(self, texts, src, tgt, **_kwargs):
        return [TRANSLATIONS[(src, tgt, text)] for text in texts]


def _source_deck(path, language: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.9))
    title_run = title.text_frame.paragraphs[0].add_run()
    title_run.text = PHRASES[language][0]
    title_run.font.bold = True
    title_run.font.size = Pt(30)
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.5), Inches(1.2))
    body_run = body.text_frame.paragraphs[0].add_run()
    body_run.text = PHRASES[language][1]
    body_run.font.size = Pt(20)
    prs.save(path)


@pytest.mark.parametrize(
    ("value", "expected"),
    [("en-US", "en"), ("DE-de", "de"), ("zh-Hans", "zh")],
)
def test_language_aliases(value: str, expected: str) -> None:
    assert normalize_language(value) == expected


def test_language_engine_metadata() -> None:
    assert language_name("zh") == "Simplified Chinese (zh-CN)"
    assert pptx_language_tag("de") == "de-DE"
    assert tesseract_language("zh") == "chi_sim+eng"
    assert tesseract_language("de") == "deu+eng"
    assert paddle_language("zh") == "ch"
    assert paddle_language("de") == "german"
    with pytest.raises(ValueError, match="supported languages"):
        normalize_language("fr")


@pytest.mark.parametrize(
    ("source", "target"),
    [
        ("en", "zh"),
        ("en", "de"),
        ("zh", "en"),
        ("zh", "de"),
        ("de", "en"),
        ("de", "zh"),
    ],
)
def test_pptx_roundtrip_all_language_directions(
    source: str, target: str, tmp_path, monkeypatch
) -> None:
    input_path = tmp_path / f"source-{source}.pptx"
    output_path = tmp_path / f"output-{source}-{target}.pptx"
    _source_deck(input_path, source)

    ir = parse_document(str(input_path), source, target, asset_dir=str(tmp_path / "assets"))
    progress_events = []
    with monkeypatch.context() as patch:
        patch.setattr(
            "translayer.localize.text_pipeline.registry.get",
            lambda kind, key: _MatrixTranslator(),
        )
        localize_text(
            ir,
            engine_name="matrix",
            progress_callback=progress_events.append,
        )
    assert progress_events[0]["completed"] == 0
    assert progress_events[-1]["completed"] == progress_events[-1]["total"]
    assert progress_events[-1]["completed_items"] == progress_events[-1]["total_items"]
    render_document(ir, str(input_path), str(output_path))

    rendered = Presentation(output_path)
    text = "\n".join(
        shape.text_frame.text
        for slide in rendered.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert PHRASES[target][0] in text
    assert PHRASES[target][1] in text
    assert rendered.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.bold
    assert rendered.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size == Pt(30)

    with zipfile.ZipFile(output_path) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert f'lang="{pptx_language_tag(target)}"' in slide_xml
    if target == "zh":
        assert 'typeface="Microsoft YaHei"' in slide_xml
