from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from translayer.localize.text_pipeline import localize_text
from translayer.parsers.pptx_parser import PptxParser
from translayer.renderers.pptx_renderer import PptxRenderer


class _BilingualTranslator:
    def translate(self, texts, **_kwargs):
        translations = {
            "谢谢": "Thank you",
            "Thank": "Thank",
            " you": " you",
        }
        return [translations[text] for text in texts]


def test_bilingual_slide_keeps_existing_target_and_removes_source(
    tmp_path,
    monkeypatch,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    source_shape = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(2),
        Inches(1),
    )
    source_shape.text = "谢谢"
    target_shape = slide.shapes.add_textbox(
        Inches(3.2),
        Inches(2),
        Inches(2),
        Inches(1),
    )
    target_shape.text_frame.paragraphs[0].text = "Thank"
    target_shape.text_frame.add_paragraph().text = " you"
    source = tmp_path / "bilingual-source.pptx"
    prs.save(source)

    parser = PptxParser()
    ir = parser.parse(
        str(source),
        {
            "source_lang": "zh",
            "target_lang": "en",
            "asset_dir": str(tmp_path / "assets"),
        },
    )
    monkeypatch.setattr(
        "translayer.localize.text_pipeline.registry.get",
        lambda *_args, **_kwargs: _BilingualTranslator(),
    )

    localize_text(ir, engine_name="bilingual")

    chinese = next(block for block in ir.blocks if block.source_text == "谢谢")
    existing_english = [block for block in ir.blocks if block.source_text != "谢谢"]
    assert chinese.target_text == ""
    assert [block.target_text for block in existing_english] == ["Thank", " you"]

    output = tmp_path / "bilingual-output.pptx"
    PptxRenderer().render(ir, str(source), str(output))
    rendered = Presentation(output)
    assert rendered.slides[0].shapes[0].text == ""
    assert rendered.slides[0].shapes[1].text == "Thank\n you"
