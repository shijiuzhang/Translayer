"""Parser + renderer round-trip and write-back tests."""

from __future__ import annotations

import hashlib

from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from translayer.parsers.pptx_parser import PptxParser
from translayer.renderers.pptx_renderer import PptxRenderer


def test_parse_extracts_blocks_and_images(sample_pptx):
    ir = PptxParser().parse(sample_pptx, {"source_lang": "en", "target_lang": "zh"})
    texts = [b.source_text for b in ir.blocks if b.type != "image"]
    assert "Quarterly Results" in texts
    assert "Revenue grew strongly." in texts
    assert "Costs were controlled." in texts
    # table cells
    assert "Metric" in texts and "Users" in texts and "1000" in texts
    # one image extracted
    assert len(ir.resources.images) == 1
    img = ir.resources.images[0]
    assert img.width == 400 and img.height == 200
    # title role detected
    assert any(b.semantic_role == "title" for b in ir.blocks)


def test_source_refs_are_addressable(sample_pptx):
    ir = PptxParser().parse(sample_pptx, {})
    for b in ir.blocks:
        if b.type == "table_cell":
            assert b.source_ref.row is not None and b.source_ref.col is not None
        elif b.type == "image":
            assert b.source_ref.image_id is not None
        else:
            assert b.source_ref.paragraph_index is not None


def test_render_writes_back_translations(sample_pptx, tmp_path):
    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(sample_pptx, {})

    mapping = {
        "Quarterly Results": "季度业绩",
        "Revenue grew strongly.": "收入大幅增长。",
        "Costs were controlled.": "成本得到控制。",
        "Metric": "指标",
        "Value": "数值",
        "Users": "用户",
        "1000": "1000",
    }
    for b in ir.blocks:
        if b.source_text in mapping:
            b.target_text = mapping[b.source_text]

    out = str(tmp_path / "out.pptx")
    renderer.render(ir, sample_pptx, out)

    prs = Presentation(out)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text_frame.text)
    blob = "\n".join(all_text)
    assert "季度业绩" in blob
    assert "收入大幅增长。" in blob
    assert "成本得到控制。" in blob
    assert "指标" in blob and "用户" in blob
    # English originals are gone
    assert "Quarterly Results" not in blob
    assert "Revenue grew strongly." not in blob


def test_roundtrip_preserves_structure(sample_pptx, tmp_path):
    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(sample_pptx, {})
    out = str(tmp_path / "identity.pptx")
    renderer.render(ir, sample_pptx, out)  # no target_text set -> identity-ish
    src, dst = Presentation(sample_pptx), Presentation(out)
    assert len(src.slides) == len(dst.slides)
    for s1, s2 in zip(src.slides, dst.slides, strict=False):
        assert len(s1.shapes) == len(s2.shapes)


def test_render_replaces_localized_image(sample_pptx, tmp_path):
    """Localized image resources must be written back into the output pptx."""
    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(sample_pptx, {})

    assert len(ir.resources.images) == 1
    image = ir.resources.images[0]

    # Create a visibly different "localized" image.
    localized_path = str(tmp_path / "localized.png")
    from PIL import Image as PILImage
    from PIL import ImageDraw

    pil = PILImage.new("RGB", (image.width, image.height), "red")
    ImageDraw.Draw(pil).text((10, 10), "localized", fill="white")
    pil.save(localized_path)
    image.localized_data_ref = localized_path

    out = str(tmp_path / "image_out.pptx")
    renderer.render(ir, sample_pptx, out)

    # Verify the picture in the output has the new blob.
    prs = Presentation(out)
    pictures = [
        shape for slide in prs.slides for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == 1
    with open(localized_path, "rb") as fh:
        expected_hash = hashlib.md5(fh.read()).hexdigest()
    assert hashlib.md5(pictures[0].image.blob).hexdigest() == expected_hash


def test_cropped_picture_is_ocrd_as_visible_area_and_crop_is_reset(tmp_path):
    source_image = tmp_path / "wide.png"
    PILImage.new("RGB", (200, 100), "white").save(source_image)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(source_image), Inches(1), Inches(1), Inches(2), Inches(1)
    )
    picture.crop_right = 0.25
    source_pptx = tmp_path / "cropped.pptx"
    prs.save(source_pptx)

    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(
        str(source_pptx),
        {"asset_dir": str(tmp_path / "assets")},
    )
    image = ir.resources.images[0]
    assert image.reset_crop_on_render
    assert (image.width, image.height) == (150, 100)

    localized = tmp_path / "localized.png"
    PILImage.new("RGB", (150, 100), "red").save(localized)
    image.localized_data_ref = str(localized)
    output = tmp_path / "output.pptx"
    renderer.render(ir, str(source_pptx), str(output))

    output_picture = next(
        shape
        for shape in Presentation(output).slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert output_picture.crop_left == 0
    assert output_picture.crop_top == 0
    assert output_picture.crop_right == 0
    assert output_picture.crop_bottom == 0


def test_renderer_keeps_complete_words_and_fits_them_into_text_box(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(1.5), Inches(0.8)
    )
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "谢谢"
    run.font.size = Pt(64)
    source = tmp_path / "fit-source.pptx"
    prs.save(source)

    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(str(source), {"source_lang": "zh", "target_lang": "en"})
    block = next(block for block in ir.blocks if block.source_text == "谢谢")
    block.target_text = "Thank you"
    output = tmp_path / "fit-output.pptx"
    renderer.render(ir, str(source), str(output))

    output_shape = Presentation(output).slides[0].shapes[0]
    output_run = output_shape.text_frame.paragraphs[0].runs[0]
    assert output_shape.text == "Thank you"
    assert output_run.font.size is not None
    assert output_run.font.size.pt < 64


def test_renderer_fits_text_before_foreground_picture(tmp_path):
    picture_path = tmp_path / "foreground.png"
    PILImage.new("RGB", (100, 100), "white").save(picture_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    text_shape = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(1)
    )
    text_shape.text = "source"
    foreground = slide.shapes.add_picture(
        str(picture_path), Inches(6), Inches(1), Inches(2), Inches(2)
    )

    visible_width = PptxRenderer._visible_text_width(
        text_shape,
        [foreground],
    )

    assert visible_width == Inches(5)

    source = tmp_path / "occluded-source.pptx"
    prs.save(source)
    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(
        str(source),
        {
            "source_lang": "en",
            "target_lang": "de",
            "asset_dir": str(tmp_path / "occluded-assets"),
        },
    )
    next(block for block in ir.blocks if block.source_text == "source").target_text = (
        "Complete translated sentence that must remain visible."
    )
    output = tmp_path / "occluded-output.pptx"
    renderer.render(ir, str(source), str(output))

    output_text_shape = Presentation(output).slides[0].shapes[0]
    assert output_text_shape.width == Inches(5)
    assert output_text_shape.text.endswith("visible.")


def test_renderer_fits_text_inside_background_container(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    container = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(1),
        Inches(6.5),
        Inches(1.2),
    )
    container.fill.solid()
    text_shape = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(1.2),
        Inches(8),
        Inches(0.8),
    )
    text_shape.text = "source"
    visible_width = PptxRenderer._container_text_width(
        text_shape,
        list(slide.shapes)[:1],
    )

    assert visible_width == Inches(6.1)

    source = tmp_path / "container-source.pptx"
    prs.save(source)
    parser, renderer = PptxParser(), PptxRenderer()
    ir = parser.parse(
        str(source),
        {
            "source_lang": "en",
            "target_lang": "de",
            "asset_dir": str(tmp_path / "container-assets"),
        },
    )
    next(block for block in ir.blocks if block.source_text == "source").target_text = (
        "Complete translated sentence that remains inside its colored container."
    )
    output = tmp_path / "container-output.pptx"
    renderer.render(ir, str(source), str(output))

    output_text_shape = Presentation(output).slides[0].shapes[1]
    assert output_text_shape.width == Inches(6.1)
    assert output_text_shape.text.endswith("container.")
