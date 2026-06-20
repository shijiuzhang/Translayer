"""End-to-end pipeline test using the offline mock engines (no network)."""

from __future__ import annotations

from pptx import Presentation

from translayer.localize.orchestrator import localize
from translayer.pipeline import (
    enrich_document,
    parse_document,
    render_document,
    translate_document,
)


def test_full_pipeline_offline(sample_pptx, tmp_path):
    out = str(tmp_path / "translated.pptx")
    ir = translate_document(
        sample_pptx, out,
        source_lang="en", target_lang="zh",
        translation_engine="mock",
        ocr_engine="mock",
        inpaint_engine="pillow",
        images=True,
    )
    # Every translatable block got a translation.
    assert all(b.target_text is not None for b in ir.translatable_blocks())
    # Length constraints were estimated for boxed blocks.
    assert any(b.constraints.max_chars for b in ir.translatable_blocks())
    # Image was localized.
    assert ir.resources.images[0].localized_data_ref is not None

    # Output is a valid pptx with the mock-translated text written back.
    prs = Presentation(out)
    assert len(prs.slides) == len(Presentation(sample_pptx).slides)


def test_pipeline_stages_separately(sample_pptx):
    ir = parse_document(sample_pptx, "en", "zh")
    enrich_document(ir, images=True, ocr_engine="mock")
    # OCR ran on the image.
    assert ir.resources.images[0].text_regions
    localize(ir, translation_engine="mock", inpaint_engine="pillow", images=True)
    assert ir.resources.images[0].localized_data_ref


def test_render_only(sample_pptx, tmp_path):
    ir = parse_document(sample_pptx, "en", "zh")
    for b in ir.translatable_blocks():
        b.target_text = "X"
    out = str(tmp_path / "r.pptx")
    render_document(ir, sample_pptx, out)
    assert Presentation(out)
