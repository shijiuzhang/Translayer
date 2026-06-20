"""Offline demo: build a sample English deck (with an in-image label), then
localize it to Chinese using the deterministic ``mock`` engines — no network,
no API keys. Produces ``output_zh.pptx``.

Run:
    python examples/demo_en_to_zh/run_demo.py
"""

from __future__ import annotations

import os

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from translayer.pipeline import translate_document

HERE = os.path.dirname(os.path.abspath(__file__))


def build_sample(path: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Business Review"
    body = slide.placeholders[1]
    body.text_frame.text = "Revenue grew 24% year over year."
    body.text_frame.add_paragraph().text = "Operating costs stayed flat."

    chart = os.path.join(HERE, "chart.png")
    img = Image.new("RGB", (480, 220), (235, 240, 250))
    ImageDraw.Draw(img).text((24, 90), "Revenue", fill=(11, 105, 218))
    img.save(chart)
    slide.shapes.add_picture(chart, Inches(5), Inches(2.2), Inches(3.5), Inches(1.6))

    prs.save(path)


def main() -> None:
    sample = os.path.join(HERE, "sample_en.pptx")
    output = os.path.join(HERE, "output_zh.pptx")
    build_sample(sample)

    translate_document(
        sample, output,
        source_lang="en", target_lang="zh",
        translation_engine="mock", ocr_engine="mock", inpaint_engine="pillow",
        images=True,
    )
    print(f"Wrote {output}")
    print("Open it to see the layout intact and text translated.")


if __name__ == "__main__":
    main()
